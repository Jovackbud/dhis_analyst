"""PPTX generator — python-pptx with health-appropriate colour palette.

Features:
- 16:9 widescreen (13.333 × 7.5 inches)
- Dark navy/teal health theme
- Title, content, chart, table slide types
- Footer with generation timestamp
"""
from __future__ import annotations

import io
import logging
from datetime import datetime, timezone
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger("dhis2_analyst.pptx_gen")


def slides_to_pptx(
    slides: list[dict[str, Any]],
    title: str = "DHIS2 Analyst Briefing",
) -> bytes:
    """Convert a slide manifest to a styled PPTX."""
    logger.info("pptx_generate_start", extra={"slide_count": len(slides), "title_len": len(title)})
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    now = datetime.now(timezone.utc).strftime("%d %B %Y %H:%M UTC")

    # --- Title slide ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    _style_text(title_slide.shapes.title, size=36, bold=True, color=RGBColor(0xF8, 0xFA, 0xFC))
    if len(title_slide.placeholders) > 1:
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"DHIS2 Public Health Intelligence Assistant\n{now}"
        _style_text(subtitle, size=16, color=RGBColor(0x9F, 0xB2, 0xBF))

    _set_slide_bg(title_slide, RGBColor(0x0B, 0x17, 0x1D))

    # --- Content slides ---
    for item in slides:
        slide_type = item.get("type", "content")

        if slide_type == "title":
            continue  # Already handled above

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        _set_slide_bg(slide, RGBColor(0x0B, 0x17, 0x1D))

        # Title
        slide.shapes.title.text = str(item.get("title", "Insight"))
        _style_text(slide.shapes.title, size=28, bold=True, color=RGBColor(0xF8, 0xFA, 0xFC))

        # Body content
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()

            content = item.get("content", "")
            if isinstance(content, list):
                content = "\n".join(str(x) for x in content)
            elif isinstance(content, dict):
                content = str(content.get("title", ""))

            for line in str(content).split("\n"):
                p = body.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
                p.font.name = "Arial"

        # Footer
        _add_footer(slide, now, prs)

    out = io.BytesIO()
    prs.save(out)
    data = out.getvalue()
    logger.info("pptx_generate_ok", extra={"size_bytes": len(data), "slide_count": len(prs.slides)})
    return data


def _style_text(shape, size: int = 14, bold: bool = False, color=None) -> None:
    """Apply font styling to all text in a shape."""
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = "Arial"
            if color:
                run.font.color.rgb = color


def _set_slide_bg(slide, color: RGBColor) -> None:
    """Set slide background to a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_footer(slide, text: str, prs: Presentation) -> None:
    """Add a small footer text box to the bottom of a slide."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(0.5),
        prs.slide_height - Inches(0.5),
        Inches(5),
        Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Generated {text} | DHIS2 AI Analyst"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x6B, 0x7B, 0x8D)
    p.font.name = "Arial"
